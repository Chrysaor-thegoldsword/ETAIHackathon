"""Utility script to create the pitch deck for ET Money Mentor.

This script uses python‑pptx to assemble a concise slide deck summarizing
the problem, solution, architecture, demo flow and impact model.  Run
this script after generating the architecture diagram.  The resulting
PPTX file can be played back locally or converted to a video using
screen capture tools.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def create_slides(diagram_path: str, decorative_image: str, impact_model_path: str, output_pptx: str) -> None:
    prs = Presentation()

    # Define colours
    et_red = RGBColor(0xA3, 0x00, 0x00)
    et_gold = RGBColor(0xCC, 0x99, 0x33)
    dark_grey = RGBColor(0x22, 0x22, 0x22)

    def add_title_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "ET Money Mentor"
        title.text_frame.paragraphs[0].font.color.rgb = et_red
        title.text_frame.paragraphs[0].font.size = Pt(40)
        subtitle.text = "Personal Finance Copilot"
        subtitle.text_frame.paragraphs[0].font.color.rgb = dark_grey
        subtitle.text_frame.paragraphs[0].font.size = Pt(20)
        # decorative image on title slide
        if os.path.isfile(decorative_image):
            slide.shapes.add_picture(decorative_image, Inches(5.5), Inches(0.5), width=Inches(4), height=Inches(3))

    def add_problem_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        # Title
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8.5), Inches(0.8))
        title_tf = title_shape.text_frame
        title_tf.text = "Problem Statement"
        title_tf.paragraphs[0].font.size = Pt(28)
        title_tf.paragraphs[0].font.bold = True
        title_tf.paragraphs[0].font.color.rgb = et_red
        # Body text
        body_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(5.5), Inches(3.0))
        body_tf = body_shape.text_frame
        body_tf.text = "Millions of Indian retail investors lack structured financial guidance."
        p = body_tf.add_paragraph()
        p.text = "They rely on tips, struggle with tax rules and overlook long‑term planning."
        p.level = 0
        p = body_tf.add_paragraph()
        p.text = "Manual research and advisory services are time‑consuming and expensive (₹25k/year)."
        p.level = 0
        # Decorative image
        if os.path.isfile(decorative_image):
            slide.shapes.add_picture(decorative_image, Inches(6.5), Inches(1.2), width=Inches(3), height=Inches(2.5))

    def add_solution_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8.5), Inches(0.8))
        title_tf = title_shape.text_frame
        title_tf.text = "Solution Overview"
        title_tf.paragraphs[0].font.size = Pt(28)
        title_tf.paragraphs[0].font.bold = True
        title_tf.paragraphs[0].font.color.rgb = et_red
        body_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(8.0), Inches(4.5))
        body_tf = body_shape.text_frame
        body_tf.text = "ET Money Mentor"
        p = body_tf.add_paragraph()
        p.text = "• Conversational onboarding captures income, expenses, assets, liabilities and goals"
        p = body_tf.add_paragraph()
        p.text = "• Multi‑agent architecture computes health score, plans goals, optimizes tax and suggests portfolios"
        p = body_tf.add_paragraph()
        p.text = "• Personalized report with actionable next steps"

    def add_architecture_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8.5), Inches(0.8))
        title_tf = title_shape.text_frame
        title_tf.text = "Architecture"
        title_tf.paragraphs[0].font.size = Pt(28)
        title_tf.paragraphs[0].font.bold = True
        title_tf.paragraphs[0].font.color.rgb = et_red
        body_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(4.8), Inches(3.5))
        body_tf = body_shape.text_frame
        body_tf.text = "Agents communicate via the orchestrator"
        p = body_tf.add_paragraph()
        p.text = "• Conversation Agent orchestrates data collection and coordination"
        p = body_tf.add_paragraph()
        p.text = "• Analysis agents (score, goals, tax, portfolio) run independently on the profile"
        p = body_tf.add_paragraph()
        p.text = "• Report generator compiles results for the user"
        # Insert diagram
        if os.path.isfile(diagram_path):
            slide.shapes.add_picture(diagram_path, Inches(5.5), Inches(1.4), width=Inches(4), height=Inches(3))

    def add_demo_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8.5), Inches(0.8))
        title_tf = title_shape.text_frame
        title_tf.text = "Demo Flow"
        title_tf.paragraphs[0].font.size = Pt(28)
        title_tf.paragraphs[0].font.bold = True
        title_tf.paragraphs[0].font.color.rgb = et_red
        body_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(8.5), Inches(4.5))
        body_tf = body_shape.text_frame
        body_tf.text = "1. User begins chat; agent asks questions"
        p = body_tf.add_paragraph()
        p.text = "2. Upload optional documents to pre‑fill fields"
        p = body_tf.add_paragraph()
        p.text = "3. System computes score and plans goals"
        p = body_tf.add_paragraph()
        p.text = "4. Tax regime compared and portfolio suggested"
        p = body_tf.add_paragraph()
        p.text = "5. User downloads personalized report"

    def add_impact_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8.5), Inches(0.8))
        title_tf = title_shape.text_frame
        title_tf.text = "Impact & Business Model"
        title_tf.paragraphs[0].font.size = Pt(28)
        title_tf.paragraphs[0].font.bold = True
        title_tf.paragraphs[0].font.color.rgb = et_red
        body_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(8.5), Inches(4.5))
        body_tf = body_shape.text_frame
        body_tf.text = "Revenue: 200 conversions per month × ₹1,000 commission = ₹2 Lakh/month"
        p = body_tf.add_paragraph()
        p.text = "User savings: up to ₹5 Crore/year in avoided advisor fees"
        p = body_tf.add_paragraph()
        p.text = "Time savings: 24,000 hours saved/month (10 min vs 5 hours)"
        p = body_tf.add_paragraph()
        p.text = "Loyalty: drives engagement across ET Wealth, Markets and Prime"

    add_title_slide()
    add_problem_slide()
    add_solution_slide()
    add_architecture_slide()
    add_demo_slide()
    add_impact_slide()

    prs.save(output_pptx)


if __name__ == '__main__':
    # Paths relative to the project root
    diagram = os.path.join('et_money_mentor', 'architecture_diagram.png')
    deco_img = os.path.join('et_money_mentor', 'decorative.png')
    impact = os.path.join('et_money_mentor', 'impact_model.md')
    output = os.path.join('et_money_mentor', 'pitch_deck.pptx')
    # Copy decorative image if not present (placeholder)
    if not os.path.isfile(deco_img):
        # The decorative image must be placed manually or generated beforehand
        pass
    create_slides(diagram, deco_img, impact, output)
    print(f"Created {output}")